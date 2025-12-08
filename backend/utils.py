import io
import re
from typing import List
from fastapi import UploadFile
from docx import Document
from pptx import Presentation
import pdfplumber
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
import html

async def read_file_to_text(file: UploadFile) -> str:
    """
    Read file bytes and return cleaned text. Handles txt/docx/pdf/pptx.
    For PDFs we use per-page extraction and then normalize.
    """
    contents = await file.read()
    name = file.filename.lower()

    if name.endswith(".txt"):
        raw = contents.decode("utf-8", errors="ignore")
        return normalize_whitespace(raw)

    if name.endswith(".docx"):
        doc = Document(io.BytesIO(contents))
        raw = "\n".join(p.text for p in doc.paragraphs)
        return normalize_whitespace(raw)

    if name.endswith(".pptx"):
        presentation = Presentation(io.BytesIO(contents))
        pieces = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pieces.append(shape.text)
        raw = "\n".join(pieces)
        return normalize_whitespace(raw)

    if name.endswith(".pdf"):
        pages_text = []
        with pdfplumber.open(io.BytesIO(contents)) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                # split into lines to preserve page structure
                lines = [ln.rstrip() for ln in text.splitlines()]
                pages_text.append(lines)

        if not any(pages_text):
            return ""  # no selectable text

        # remove repeated headers/footers across pages
        pages_text = remove_repeated_header_footer(pages_text)

        # join lines inside pages with heuristics
        page_paragraphs = []
        for lines in pages_text:
            joined = join_broken_lines(lines)
            page_paragraphs.append(joined)

        # join pages with a page-break marker (double newline)
        raw = "\n\n".join(page_paragraphs)
        raw = normalize_whitespace(raw)
        return raw

    # default fallback
    return contents.decode("utf-8", errors="ignore")


# -------------------------
# Text cleaning helpers
# -------------------------
def normalize_whitespace(text: str) -> str:
    """Normalize whitespace, fix many spaces/newlines, convert weird non-breaking spaces."""
    if not text:
        return ""
    text = text.replace("\u00A0", " ")
    text = text.replace("\r", "")
    # collapse multiple spaces
    text = re.sub(r"[ \t]{2,}", " ", text)
    # unify many newlines into max two
    text = re.sub(r"\n{3,}", "\n\n", text)
    # trim spaces on each line
    text = "\n".join([ln.strip() for ln in text.splitlines()])
    # remove leading/trailing
    return text.strip()


def remove_repeated_header_footer(pages_lines: List[List[str]]) -> List[List[str]]:
    """
    Detect lines that repeat at top or bottom of many pages and remove them.
    Strategy:
      - Collect top N and bottom N lines for each page (non-empty).
      - Count occurrences across pages. If a line appears on >50% pages at same
        relative position, treat as header/footer and remove.
    """
    if not pages_lines:
        return pages_lines

    top_n = 3
    bottom_n = 3
    top_counts = {}
    bottom_counts = {}
    page_count = len(pages_lines)

    # collect candidates
    for lines in pages_lines:
        # normalized small-lines list
        nonempty = [ln.strip() for ln in lines if ln.strip()]
        # top candidates
        for i in range(min(top_n, len(nonempty))):
            key = normalize_short(nonempty[i])
            top_counts[key] = top_counts.get(key, 0) + 1
        # bottom candidates
        for i in range(min(bottom_n, len(nonempty))):
            key = normalize_short(nonempty[-1 - i])
            bottom_counts[key] = bottom_counts.get(key, 0) + 1

    # decide headers/footers present on >50% pages
    header_candidates = {k for k, v in top_counts.items() if v > page_count * 0.5}
    footer_candidates = {k for k, v in bottom_counts.items() if v > page_count * 0.5}

    cleaned_pages = []
    for lines in pages_lines:
        cleaned = []
        nonempty = [ln for ln in lines]  # preserve positions
        # remove top matches
        idx = 0
        # skip over header candidates at top
        while idx < len(nonempty):
            if nonempty[idx].strip() and normalize_short(nonempty[idx]) in header_candidates:
                idx += 1
                continue
            break
        # collect middle part
        middle = nonempty[idx:]
        # remove footer candidates from bottom
        while middle and middle[-1].strip() and normalize_short(middle[-1]) in footer_candidates:
            middle.pop()
        # final cleaned lines (strip)
        cleaned = [ln for ln in (l.strip() for l in middle) if ln != ""]
        cleaned_pages.append(cleaned)

    return cleaned_pages


def normalize_short(s: str) -> str:
    """Compact a short string for header/footer comparison."""
    s = s.strip().lower()
    s = re.sub(r"\s+", " ", s)
    # remove page number patterns like 'page 1 of 5'
    s = re.sub(r"page\s*\d+(\s*of\s*\d+)?", "", s)
    s = s.strip()
    return s


def join_broken_lines(lines: List[str]) -> str:
    """
    Heuristic reflow for lines:
      - If a line ends with '-' then it's hyphenation: remove '-' and join.
      - If a line doesn't end with sentence end (.!?;:) and next line starts with lowercase/continuation,
        join with space.
      - If current line is very short and next line starts with lowercase, join.
      - Otherwise keep newline (paragraph boundary).
    """
    out_lines = []
    i = 0
    while i < len(lines):
        cur = lines[i].rstrip()
        if cur == "":
            # preserve paragraph break
            out_lines.append("")
            i += 1
            continue

        # accumulate run
        run = cur
        j = i + 1
        while j < len(lines):
            nxt = lines[j].lstrip()
            if nxt == "":
                break  # paragraph break
            # hyphenation case
            if run.endswith("-"):
                run = run[:-1] + nxt  # remove hyphen and join immediately
                j += 1
                continue
            # if current ends with punctuation, keep as new line
            if re.search(r"[\.!\?:]" + r"\s*$", run):
                break
            # if next starts with lowercase or digit, probably a continuation
            if re.match(r"^[a-z0-9]", nxt):
                run = run + " " + nxt
                j += 1
                continue
            # if current is short (<=40 chars) and next starts lowercase, join
            if len(run) < 40 and re.match(r"^[a-z]", nxt):
                run = run + " " + nxt
                j += 1
                continue
            # otherwise do not join
            break

        out_lines.append(run.strip())
        i = j
    # join runs with single newline; treat empty strings as paragraph separators
    paragraphs = []
    cur_para = []
    for ln in out_lines:
        if ln == "":
            if cur_para:
                paragraphs.append(" ".join(cur_para).strip())
                cur_para = []
        else:
            cur_para.append(ln)
    if cur_para:
        paragraphs.append(" ".join(cur_para).strip())

    # return page text with paragraphs separated by blank line
    return "\n\n".join(paragraphs)


# -------------------------
# PDF builder (robust multi-page)
# -------------------------
# 




def build_pdf_from_text_or_markdown(content: str) -> bytes:
   
   

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter,
                            topMargin=36, bottomMargin=36,
                            leftMargin=40, rightMargin=40)
    styles = getSampleStyleSheet()
    normal_style = styles["Normal"]
    story = []

    lines = content.splitlines()
    i = 0

    def chunk_paragraph(paragraph: str, max_chars: int = 900) -> List[str]:
        if len(paragraph) <= max_chars:
            return [paragraph]
        sentences = re.split(r'(?<=[.!?])\s+', paragraph)
        chunks, cur = [], ""
        for s in sentences:
            if len(cur) + len(s) + 1 <= max_chars:
                cur = (cur + " " + s).strip()
            else:
                if cur:
                    chunks.append(cur.strip())
                cur = s
        if cur:
            chunks.append(cur.strip())
        # further split very long chunks
        final = []
        for c in chunks:
            if len(c) <= max_chars:
                final.append(c)
            else:
                for j in range(0, len(c), max_chars):
                    final.append(c[j:j+max_chars])
        return final

    def is_table_start(line: str, next_line: str = "") -> bool:
        if "|" not in line:
            return False
        if re.match(r'^[\|\-\s:]+$', next_line):
            return True
        return True if "|" in line else False

    while i < len(lines):
        line = lines[i].strip()
        next_line = lines[i+1].strip() if i+1 < len(lines) else ""

        if is_table_start(line, next_line):
            # Collect table lines
            table_lines = []
            while i < len(lines) and "|" in lines[i]:
                table_lines.append(lines[i])
                i += 1

            # Skip separator line if present
            if len(table_lines) > 1 and re.match(r'^[\|\-\s:]+$', table_lines[1]):
                table_lines.pop(1)

            # Convert to 2D list for ReportLab safely
            table_data = []
            for tbl_line in table_lines:
                # get cells between pipes
                cells = [c.strip() for c in tbl_line.split("|")[1:-1]]
                if not cells or all(not c for c in cells):
                    continue  # skip empty rows
                row = [Paragraph(html.escape(c), normal_style) for c in cells]
                table_data.append(row)

            # Only build table if valid
            if table_data and len(table_data[0]) > 0:
                tbl = Table(table_data, repeatRows=1)
                tbl.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.lightgrey),
                    ('TEXTCOLOR', (0,0), (-1,0), colors.black),
                    ('ALIGN',(0,0),(-1,-1),'LEFT'),
                    ('VALIGN',(0,0),(-1,-1),'TOP'),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.black),
                    ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold')
                ]))
                story.append(tbl)
                story.append(Spacer(1, 12))
        else:
            # Collect paragraph lines (unchanged)
            para_lines = []
            while i < len(lines) and lines[i].strip() != "" and "|" not in lines[i]:
                para_lines.append(lines[i].strip())
                i += 1
            if para_lines:
                joined = []
                j = 0
                while j < len(para_lines):
                    cur = para_lines[j]
                    k = j + 1
                    while k < len(para_lines):
                        nxt = para_lines[k]
                        if cur.endswith("-"):
                            cur = cur[:-1] + nxt
                            k += 1
                            continue
                        if re.search(r"[\.!\?:]\s*$", cur):
                            break
                        if re.match(r"^[a-z0-9]", nxt) or len(cur) < 40:
                            cur = cur + " " + nxt
                            k += 1
                        else:
                            break
                    joined.append(cur.strip())
                    j = k
                for para in joined:
                    for chunk in chunk_paragraph(para):
                        safe_chunk = re.sub(r"</?para>", "", chunk, flags=re.IGNORECASE)
                        safe_chunk = html.escape(safe_chunk).replace("\n", "<br/>")
                        story.append(Paragraph(safe_chunk, normal_style))
                        story.append(Spacer(1, 8))
            else:
                i += 1

    doc.build(story)
    buffer.seek(0)
    return buffer.read()

# def build_docx_from_text(text: str) -> bytes:
#     """
#     Create a DOCX file from cleaned text with basic formatting.
#     """
#     buffer = io.BytesIO()
#     doc = Document()
#     paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]

#     for para in paragraphs:
#         doc.add_paragraph(para)
#         doc.add_paragraph("")  # add a blank line between paragraphs

#     doc.save(buffer)
#     buffer.seek(0)
#     return buffer.read()
def build_docx_from_text(content: str) -> bytes:
    import html
    buffer = io.BytesIO()
    doc = Document()

    lines = content.splitlines()
    i = 0

    # ---------- same logic as PDF ----------
    def chunk_paragraph(paragraph: str, max_chars: int = 900) -> List[str]:
        if len(paragraph) <= max_chars:
            return [paragraph]

        sentences = re.split(r"(?<=[.!?])\s+", paragraph)
        chunks, cur = [], ""

        for s in sentences:
            if len(cur) + len(s) + 1 <= max_chars:
                cur = (cur + " " + s).strip()
            else:
                if cur:
                    chunks.append(cur.strip())
                cur = s
        if cur:
            chunks.append(cur.strip())

        final = []
        for c in chunks:
            if len(c) <= max_chars:
                final.append(c)
            else:
                for j in range(0, len(c), max_chars):
                    final.append(c[j:j+max_chars])
        return final

    def is_table_start(line: str, next_line: str = "") -> bool:
        if "|" not in line:
            return False
        if re.match(r"^[\|\-\s:]+$", next_line):
            return True
        return "|" in line

    # ---------- main parsing loop ----------
    while i < len(lines):
        line = lines[i].strip()
        next_line = lines[i + 1].strip() if i + 1 < len(lines) else ""

        # ---------- TABLE HANDLING ----------
        if is_table_start(line, next_line):
            table_lines = []
            while i < len(lines) and "|" in lines[i]:
                table_lines.append(lines[i])
                i += 1

            # remove markdown separator row
            if len(table_lines) > 1 and re.match(r"^[\|\-\s:]+$", table_lines[1]):
                table_lines.pop(1)

            # convert to rows
            table_data = []
            for tbl_line in table_lines:
                cells = [c.strip() for c in tbl_line.split("|")[1:-1]]
                if not any(cells):
                    continue
                table_data.append(cells)

            if table_data:
                t = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                t.style = "Table Grid"

                for r_idx, row in enumerate(table_data):
                    for c_idx, cell in enumerate(row):
                        t.cell(r_idx, c_idx).text = cell

                doc.add_paragraph("")  # spacing after table

        else:
            # ---------- PARAGRAPH HANDLING ----------
            para_lines = []
            while i < len(lines) and lines[i].strip() != "" and "|" not in lines[i]:
                para_lines.append(lines[i].strip())
                i += 1

            if para_lines:
                # merge broken lines (same logic as PDF)
                joined = []
                j = 0
                while j < len(para_lines):
                    cur = para_lines[j]
                    k = j + 1

                    while k < len(para_lines):
                        nxt = para_lines[k]

                        if cur.endswith("-"):
                            cur = cur[:-1] + nxt
                            k += 1
                            continue

                        if re.search(r"[.!?:]\s*$", cur):
                            break

                        if re.match(r"^[a-z0-9]", nxt) or len(cur) < 40:
                            cur = cur + " " + nxt
                            k += 1
                        else:
                            break

                    joined.append(cur.strip())
                    j = k

                # chunk long paragraphs
                for para in joined:
                    for chunk in chunk_paragraph(para):
                        safe = html.escape(chunk)
                        doc.add_paragraph(safe)
                doc.add_paragraph("")  # paragraph spacing
            else:
                i += 1

    # ---------- output ----------
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()


async def extract_images_from_pdf(file: UploadFile) -> list[io.BytesIO]:
    """
    Extract all images (diagrams) from the uploaded PDF and return as list of BytesIO.
    """
    images = []
    contents = await file.read()
    doc = fitz.open(stream=contents, filetype="pdf")
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = io.BytesIO(base_image["image"])
            images.append(image_bytes)
    return images

